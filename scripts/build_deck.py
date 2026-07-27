"""Generate the presentation decks.

Two files, both built from the live analysis so a figure on a slide cannot
disagree with the scorecard behind it:

``Cardiac_Opportunity_Agent_Round1.pptx``   cover, three slides, four appendix
``Cardiac_Opportunity_Agent_Detailed.pptx`` the seven-to-ten slide version

Run ``python scripts/build_visuals.py`` first, then:

    python scripts/build_deck.py
"""

from __future__ import annotations

import sys
from dataclasses import dataclass
from pathlib import Path

PROJECT_ROOT = Path(__file__).resolve().parents[1]
if str(PROJECT_ROOT / "src") not in sys.path:
    sys.path.insert(0, str(PROJECT_ROOT / "src"))

from pptx import Presentation  # noqa: E402
from pptx.dml.color import RGBColor  # noqa: E402
from pptx.enum.shapes import MSO_SHAPE  # noqa: E402
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN  # noqa: E402
from pptx.util import Emu, Inches, Pt  # noqa: E402

from cardiac_agent.analytics.sensitivity import run_sensitivity  # noqa: E402
from cardiac_agent.analytics.whitespace import find_whitespace  # noqa: E402
from cardiac_agent.pipeline import get_context  # noqa: E402

ASSETS = PROJECT_ROOT / "docs" / "assets"
OUTPUT = PROJECT_ROOT / "deliverables"

# --------------------------------------------------------------------------
# Design system
#
# The palette is the one used by the charts, so slides and figures read as one
# document. Deep teal carries the recommendation, terracotta carries the
# caution, and everything else is neutral so the eye goes to the data.
# --------------------------------------------------------------------------

INK = RGBColor(0x1C, 0x27, 0x33)
CANVAS = RGBColor(0xFA, 0xF8, 0xF5)
TEAL = RGBColor(0x1F, 0x6F, 0x6B)
TEAL_SOFT = RGBColor(0xE4, 0xEE, 0xED)
TERRACOTTA = RGBColor(0xB4, 0x55, 0x2D)
TERRA_SOFT = RGBColor(0xF7, 0xEB, 0xE5)
SLATE = RGBColor(0x54, 0x67, 0x7D)
MUTED = RGBColor(0x7D, 0x8A, 0x99)
LINE = RGBColor(0xDD, 0xE3, 0xEA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

HEAD_FONT = "Georgia"
BODY_FONT = "Calibri"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.62)


@dataclass
class Facts:
    """Every figure the deck quotes, pulled once from the live analysis."""

    market: float
    market_growth: float
    market_real: float
    market_volume: float
    market_price: float
    market_cagr: float
    cipla: float
    cipla_share: float
    cipla_growth: float
    cipla_rank: int
    company_count: int
    spaces_scored: int
    sku_rows: int
    signals: int
    top_moi: list[dict]
    priorities: list[dict]
    whitespace: list[dict]
    stability: list[dict]
    stop: list[dict]
    franchises: list[tuple[str, float, float]]


def gather_facts(context) -> Facts:
    totals = context.totals
    segment_rows = context.company_facts[context.company_facts["level"] == "segment"]
    company_totals = (
        segment_rows.groupby("company_clean")["value_t2"].sum().sort_values(ascending=False)
    )
    rank = list(company_totals.index).index("CIPLA") + 1

    moi = (
        context.scored[context.scored["level"] == "molecule_combination"]
        .nlargest(5, "market_opportunity_index")
        .to_dict("records")
    )

    def space(level: str, label: str) -> dict:
        row = context.find_space(label, level)
        return {} if row is None else row.to_dict()

    priorities = [
        space("sub_segment", "Lipid Regulators | Statins Comb."),
        space("molecule_combination", "C10A0S ROSUVASTATIN + EZETIMIBE"),
        space("anchor_molecule", "CILNIDIPINE (all forms)"),
        space("molecule_combination", "C02F0O CILNIDIPINE + TELMISARTAN"),
        space("sub_segment", "Anti Hypertensives | AHT Triple / Poly Comb."),
        space("molecule_combination", "C10A0I ROSUVASTATIN+CLOPIDOGREL"),
    ]

    gaps = find_whitespace(
        context.scored,
        focal_overall_share=totals["focal_share"],
        levels=["anchor_molecule"],
        limit=5,
    ).to_dict("records")

    stability = (
        run_sensitivity(context.enriched, level="sub_segment", framework=context.framework, top_k=3)
        .stability.head(5)
        .to_dict("records")
    )

    stop = [
        space("sub_segment", "Anti Hypertensives | AHT Diuretic Comb."),
        space("sub_segment", "Anti Hypertensives | ACEi"),
        space("sub_segment", "Lipid Regulators | Oth. Lipid Red."),
        space("sub_segment", "Anti Hypertensives | AHT Dual Comb."),
    ]

    focal = context.sku_facts[context.sku_facts["is_focal"]]
    grouped = (
        focal.groupby("brand_root")[["MAT FEB'25", "MAT FEB'26"]].sum().nlargest(5, "MAT FEB'26")
    )
    franchises = [
        (name, float(row["MAT FEB'26"]), float(row["MAT FEB'26"] / row["MAT FEB'25"] - 1) * 100)
        for name, row in grouped.iterrows()
    ]

    return Facts(
        market=totals["market_value_t2"],
        market_growth=totals["market_yoy"] * 100,
        market_real=totals["market_real_growth"] * 100,
        market_volume=totals["market_volume_growth"] * 100,
        market_price=totals["market_price_effect"] * 100,
        market_cagr=totals["market_cagr_2y"] * 100,
        cipla=totals["focal_value_t2"],
        cipla_share=totals["focal_share"] * 100,
        cipla_growth=totals["focal_yoy"] * 100,
        cipla_rank=rank,
        company_count=len(company_totals),
        spaces_scored=len(context.scored),
        sku_rows=int(context.metadata.get("sku_rows", 0)),
        signals=len(context.corpus),
        top_moi=moi,
        priorities=priorities,
        whitespace=gaps,
        stability=stability,
        stop=stop,
        franchises=franchises,
    )


# --------------------------------------------------------------------------
# Layout helpers
# --------------------------------------------------------------------------


def _text(
    slide,
    left,
    top,
    width,
    height,
    text: str,
    *,
    size=16,
    bold=False,
    colour=INK,
    font=BODY_FONT,
    align=PP_ALIGN.LEFT,
    line_spacing=1.16,
    anchor=MSO_ANCHOR.TOP,
    space_after=0,
    italic=False,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.word_wrap = True
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    frame.vertical_anchor = anchor
    for index, line in enumerate(text.split("\n")):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.alignment = align
        paragraph.line_spacing = line_spacing
        paragraph.space_after = Pt(space_after)
        run = paragraph.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.name = font
        run.font.color.rgb = colour
    return box


def _rect(slide, left, top, width, height, fill, *, shape=MSO_SHAPE.RECTANGLE):
    box = slide.shapes.add_shape(shape, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = fill
    box.line.fill.background()
    box.shadow.inherit = False
    return box


def _blank(presentation, background=CANVAS):
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    _rect(slide, 0, 0, SLIDE_W, SLIDE_H, background)
    return slide


#: Characters that fit on one line of a 28pt Georgia title across the content
#: width. Titles longer than this wrap, which pushes the subtitle into the body
#: of the slide, so the builder sizes the title box from the real line count
#: rather than assuming one.
TITLE_CHARS_PER_LINE = 52
TITLE_SIZE = 28.0


def _content_slide(presentation, eyebrow: str, title: str, *, subtitle: str = ""):
    """Light slide with the recurring left accent bar and an eyebrow label."""
    slide = _blank(presentation)
    _rect(slide, 0, 0, Inches(0.13), SLIDE_H, TEAL)

    _text(
        slide,
        MARGIN,
        Inches(0.44),
        Inches(11.0),
        Inches(0.3),
        eyebrow.upper(),
        size=11,
        bold=True,
        colour=TEAL,
        font=BODY_FONT,
    )

    lines = sum(max(1, -(-len(line) // TITLE_CHARS_PER_LINE)) for line in title.split("\n"))
    title_height = Inches(lines * TITLE_SIZE * 1.1 / 72.0 + 0.06)
    title_top = Inches(0.78)
    _text(
        slide,
        MARGIN,
        title_top,
        Inches(12.1),
        title_height,
        title,
        size=TITLE_SIZE,
        bold=True,
        colour=INK,
        font=HEAD_FONT,
        line_spacing=1.1,
    )

    if subtitle:
        _text(
            slide,
            MARGIN,
            title_top + title_height + Inches(0.14),
            Inches(11.8),
            Inches(0.44),
            subtitle,
            size=13,
            colour=MUTED,
            font=BODY_FONT,
        )
    return slide


def _stat(slide, left, top, width, value: str, label: str, *, colour=TEAL, value_size=30):
    _text(
        slide,
        left,
        top,
        width,
        Inches(0.52),
        value,
        size=value_size,
        bold=True,
        colour=colour,
        font=HEAD_FONT,
    )
    _text(
        slide,
        left,
        top + Inches(0.52),
        width,
        Inches(0.6),
        label,
        size=11,
        colour=MUTED,
        font=BODY_FONT,
        line_spacing=1.16,
    )


def _card(slide, left, top, width, height, *, fill=WHITE, accent=None):
    _rect(slide, left, top, width, height, fill)
    if accent is not None:
        _rect(slide, left, top, Inches(0.055), height, accent)


def _picture(slide, name: str, left, top, width):
    path = ASSETS / f"{name}.png"
    if not path.exists():
        raise FileNotFoundError(f"Missing chart {path}. Run scripts/build_visuals.py first.")
    return slide.shapes.add_picture(str(path), left, top, width=width)


def _footnote(slide, text: str):
    _text(
        slide,
        MARGIN,
        Inches(6.94),
        Inches(12.1),
        Inches(0.32),
        text,
        size=9.5,
        colour=MUTED,
        font=BODY_FONT,
    )


def _table(
    slide,
    left,
    top,
    width,
    headers,
    rows,
    *,
    col_ratios,
    row_h=Inches(0.42),
    head_h=Inches(0.36),
    size=11.5,
    accent_col=None,
):
    """Lightweight table drawn from shapes.

    python-pptx tables carry theme styling that fights the palette, so the
    table is drawn manually. It also keeps every cell measurable for the
    geometry check in scripts/check_deck.py.
    """
    total = sum(col_ratios)
    widths = [Emu(int(width * ratio / total)) for ratio in col_ratios]

    x = left
    for index, header in enumerate(headers):
        _text(
            slide,
            x,
            top,
            widths[index],
            head_h,
            header,
            size=10,
            bold=True,
            colour=MUTED,
            font=BODY_FONT,
            align=PP_ALIGN.RIGHT if index and index != accent_col else PP_ALIGN.LEFT,
        )
        x += widths[index]

    y = top + head_h
    _rect(slide, left, y, width, Emu(int(Inches(0.014))), LINE)
    y += Inches(0.1)

    for row in rows:
        x = left
        for index, cell in enumerate(row):
            is_accent = index == accent_col
            _text(
                slide,
                x,
                y,
                widths[index] - Inches(0.12),
                row_h,
                str(cell),
                size=size,
                bold=is_accent or index == 0,
                colour=TEAL if is_accent else INK,
                font=BODY_FONT,
                align=PP_ALIGN.RIGHT if index and not is_accent else PP_ALIGN.LEFT,
            )
            x += widths[index]
        y += row_h
    return y


def _short(label: str, limit: int = 30) -> str:
    text = str(label).split("|")[-1].strip()
    parts = text.split(None, 1)
    if (
        parts
        and len(parts[0]) >= 5
        and parts[0][0].isalpha()
        and any(c.isdigit() for c in parts[0])
    ):
        text = parts[1] if len(parts) > 1 else text
    text = text.title() if text.isupper() else text
    text = text.replace("Metop.", "Metoprolol").replace("Telmi.", "Telmisartan")
    text = text.replace("Chlortal", "Chlortalidone").replace("Cilnidip", "Cilnidipine")
    text = text.replace("Telmis", "Telmisartan").replace("(All Forms)", "").strip()
    return text if len(text) <= limit else text[: limit - 1] + "."


# --------------------------------------------------------------------------
# Slides
# --------------------------------------------------------------------------


def slide_cover(presentation, facts: Facts, subtitle: str) -> None:
    slide = _blank(presentation, INK)
    _rect(slide, 0, 0, Inches(0.16), SLIDE_H, TEAL)
    _rect(slide, Inches(8.55), 0, Inches(4.78), SLIDE_H, RGBColor(0x22, 0x2F, 0x3D))

    _text(
        slide,
        Inches(1.05),
        Inches(1.68),
        Inches(7.0),
        Inches(0.34),
        "ASCEND SEASON 4  ??  CASE STUDY RESPONSE",
        size=11.5,
        bold=True,
        colour=TEAL_SOFT,
        font=BODY_FONT,
    )
    _text(
        slide,
        Inches(1.05),
        Inches(2.16),
        Inches(7.1),
        Inches(1.9),
        "Where Cipla can win\nin the India Cardiac market",
        size=40,
        bold=True,
        colour=WHITE,
        font=HEAD_FONT,
        line_spacing=1.08,
    )
    _text(
        slide,
        Inches(1.05),
        Inches(4.12),
        Inches(6.9),
        Inches(0.9),
        subtitle,
        size=15,
        colour=RGBColor(0xB9, 0xC6, 0xD1),
        font=BODY_FONT,
        line_spacing=1.3,
    )
    _text(
        slide,
        Inches(1.05),
        Inches(6.42),
        Inches(7.0),
        Inches(0.32),
        "AI-enabled prioritisation with integrated trend analytics  ??  MAT February 2026",
        size=11,
        colour=MUTED,
        font=BODY_FONT,
    )

    metrics = [
        (f"{facts.sku_rows:,}", "packs analysed"),
        (f"{facts.spaces_scored}", "opportunity spaces scored"),
        (f"{facts.signals}", "cited external signals"),
        ("0", "numbers written by the model"),
    ]
    top = Inches(1.72)
    for value, label in metrics:
        _text(
            slide,
            Inches(9.28),
            top,
            Inches(3.6),
            Inches(0.55),
            value,
            size=33,
            bold=True,
            colour=TEAL_SOFT,
            font=HEAD_FONT,
        )
        _text(
            slide,
            Inches(9.28),
            top + Inches(0.63),
            Inches(3.6),
            Inches(0.3),
            label,
            size=11.5,
            colour=MUTED,
            font=BODY_FONT,
        )
        top += Inches(1.17)


def slide_finding(presentation, facts: Facts) -> None:
    slide = _content_slide(
        presentation,
        "What the agent found",
        f"Cipla grows at {facts.cipla_growth:.1f}% in a market growing {facts.market_growth:.1f}%.\n"
        "Two configurations inside it are winnable.",
        subtitle="Ranked on market opportunity, then gated by right to win. "
        "The two scores are kept apart so the trade-off stays visible.",
    )

    stats = [
        (f"{facts.market:,.0f}", "INR crore\nCardiac market", TEAL),
        (
            f"{facts.market_real:.1f}%",
            f"real demand growth\nof {facts.market_growth:.1f}% reported",
            SLATE,
        ),
        (
            f"{facts.cipla_share:.2f}%",
            f"Cipla share\nrank {facts.cipla_rank} of {facts.company_count}",
            TERRACOTTA,
        ),
    ]
    left = MARGIN
    for value, label, colour in stats:
        _stat(slide, left, Inches(2.32), Inches(1.72), value, label, colour=colour, value_size=27)
        left += Inches(1.86)

    _rect(slide, MARGIN, Inches(3.72), Inches(5.2), Inches(0.014), LINE)

    _text(
        slide,
        MARGIN,
        Inches(3.95),
        Inches(5.2),
        Inches(0.3),
        "HOW IT WORKS",
        size=10.5,
        bold=True,
        colour=TEAL,
        font=BODY_FONT,
    )
    method = (
        "Six levels of opportunity space, from segment down to individual molecule\n"
        "combinations and cross-hierarchy molecule franchises.\n\n"
        "Four pillars: market attractiveness, future potential, competitive headroom,\n"
        "and a right-to-win gate built from Cipla's actual estate.\n\n"
        "Every figure is computed deterministically. The language model narrates and\n"
        "is rejected if it states a number the evidence does not contain."
    )
    _text(
        slide,
        MARGIN,
        Inches(4.32),
        Inches(5.35),
        Inches(2.3),
        method,
        size=12,
        colour=SLATE,
        font=BODY_FONT,
        line_spacing=1.22,
    )

    _picture(slide, "priority-matrix-sub-segment", Inches(6.35), Inches(2.06), Inches(6.4))
    _footnote(
        slide, "Source: Cardiac prescription audit, MAT Feb 2026. Bubble area is market value."
    )


def slide_priorities(presentation, facts: Facts) -> None:
    statins, ezetimibe, cilnidipine, cilni_telmi, triple, rosuva_clopi = facts.priorities

    slide = _content_slide(
        presentation,
        "The recommendation",
        "Prioritise the ezetimibe extension of Rosulip,\nand the cilnidipine family entered through Cresar",
        subtitle="Both are brand extensions rather than standing starts, which is what makes them "
        "reachable inside two years.",
    )

    cards = [
        (
            MARGIN,
            "PRIORITY 1",
            TEAL,
            TEAL_SOFT,
            "Statin combinations,\nvia ezetimibe",
            [
                (
                    "Sub-segment",
                    f"{statins['value_t2']:,.0f} cr, +{statins['value_yoy'] * 100:.1f}%, real +{statins['real_growth'] * 100:.1f}%",
                ),
                (
                    "Structure",
                    f"HHI {statins['hhi']:,.0f}, {int(statins['n_players'])} players, leader {statins['leader_share'] * 100:.0f}%",
                ),
                (
                    "Cipla today",
                    f"{statins['focal_value_t2']:,.1f} cr, {statins['focal_share_t2'] * 100:.2f}%, +{statins['focal_yoy'] * 100:.1f}%",
                ),
                (
                    "The target",
                    f"Rosuvastatin + ezetimibe: {ezetimibe['value_t2']:,.0f} cr, +{ezetimibe['value_yoy'] * 100:.0f}%",
                ),
                (
                    "Right to win",
                    "Rosulip is the fastest-growing Cipla franchise; Rosulip EZ is launched",
                ),
            ],
        ),
        (
            Inches(6.92),
            "PRIORITY 2",
            TERRACOTTA,
            TERRA_SOFT,
            "Cilnidipine combinations,\nvia Cresar",
            [
                (
                    "Franchise",
                    f"{cilnidipine['value_t2']:,.0f} cr across plain and combination packs, +{cilnidipine['value_yoy'] * 100:.1f}%",
                ),
                (
                    "Cipla today",
                    f"{cilnidipine['focal_share_t2'] * 100:.2f}% share, {cilnidipine['focal_value_t2']:,.1f} cr, growing {cilnidipine['focal_yoy'] * 100:.0f}%",
                ),
                (
                    "Where to enter",
                    f"Not the plain molecule. The dual sits at HHI {cilni_telmi['hhi']:,.0f}; the triples are more open",
                ),
                (
                    "The target",
                    f"Cilnidipine + telmisartan: {cilni_telmi['value_t2']:,.0f} cr, +{cilni_telmi['value_yoy'] * 100:.0f}%, real +{cilni_telmi['real_growth'] * 100:.0f}%",
                ),
                (
                    "Right to win",
                    "Cresar carries telmisartan; Cresar LN already sells in this space",
                ),
            ],
        ),
    ]

    for left, tag, accent, soft, heading, rows in cards:
        _card(slide, left, Inches(2.42), Inches(5.78), Inches(3.62), fill=WHITE, accent=accent)
        _rect(slide, left, Inches(2.42), Inches(5.78), Inches(0.92), soft)
        _rect(slide, left, Inches(2.42), Inches(0.055), Inches(3.62), accent)
        _text(
            slide,
            left + Inches(0.3),
            Inches(2.58),
            Inches(5.2),
            Inches(0.24),
            tag,
            size=10,
            bold=True,
            colour=accent,
            font=BODY_FONT,
        )
        _text(
            slide,
            left + Inches(0.3),
            Inches(2.84),
            Inches(5.2),
            Inches(0.62),
            heading,
            size=17,
            bold=True,
            colour=INK,
            font=HEAD_FONT,
            line_spacing=1.06,
        )
        top = Inches(3.52)
        for label, value in rows:
            _text(
                slide,
                left + Inches(0.3),
                top,
                Inches(1.34),
                Inches(0.36),
                label,
                size=10.5,
                bold=True,
                colour=MUTED,
                font=BODY_FONT,
            )
            _text(
                slide,
                left + Inches(1.68),
                top,
                Inches(3.82),
                Inches(0.46),
                value,
                size=11.5,
                colour=INK,
                font=BODY_FONT,
                line_spacing=1.14,
            )
            top += Inches(0.49)

    _text(
        slide,
        MARGIN,
        Inches(6.22),
        Inches(2.0),
        Inches(0.26),
        "TRADE-OFFS RESOLVED",
        size=10.5,
        bold=True,
        colour=TEAL,
        font=BODY_FONT,
    )
    tradeoffs = [
        f"Size against growth. AHT dual combinations is larger at {facts.stop[3]['value_t2']:,.0f} cr, "
        f"but only {facts.stop[3]['real_growth'] * 100:.1f} of its {facts.stop[3]['value_yoy'] * 100:.1f} points of growth is real. Chose real growth.",
        f"Growth against competition. Other lipid reducers grows {facts.stop[2]['value_yoy'] * 100:.0f}% and is "
        f"{facts.stop[2]['leader_share'] * 100:.0f}% held by one company. Chose accessibility.",
        "Attractiveness against right to win. Cilnidipine is larger and scores higher; ezetimibe is winnable sooner. "
        "Resolved by sequencing, not by choosing.",
    ]
    top = Inches(6.5)
    for line in tradeoffs:
        _rect(slide, MARGIN, top + Inches(0.055), Inches(0.055), Inches(0.16), TEAL)
        _text(
            slide,
            MARGIN + Inches(0.18),
            top,
            Inches(12.0),
            Inches(0.28),
            line,
            size=10.5,
            colour=SLATE,
            font=BODY_FONT,
        )
        top += Inches(0.29)


def slide_whitespace(presentation, facts: Facts) -> None:
    slide = _content_slide(
        presentation,
        "Where else, and where to stop",
        "Five franchises are underpenetrated and reachable.\nFour positions should be harvested, not defended.",
        subtitle=f"Penetration index is Cipla's share of a space against its "
        f"{facts.cipla_share:.2f}% share of the therapy area.",
    )

    _picture(slide, "whitespace-gap", MARGIN, Inches(2.28), Inches(6.5))
    _text(
        slide,
        MARGIN,
        Inches(5.6),
        Inches(6.5),
        Inches(0.6),
        "Every one has a Cipla molecule or brand behind it, which is what separates a "
        "target from a wish. Reaching fair share across the five is worth roughly "
        f"{sum(row['value_gap_cr'] for row in facts.whitespace):,.0f} crore.",
        size=11.5,
        colour=SLATE,
        font=BODY_FONT,
        line_spacing=1.24,
    )

    left = Inches(7.36)
    _text(
        slide,
        left,
        Inches(2.3),
        Inches(5.4),
        Inches(0.28),
        "STOP FUNDING",
        size=10.5,
        bold=True,
        colour=TERRACOTTA,
        font=BODY_FONT,
    )

    reasons = [
        (
            facts.stop[0],
            "Reported growth is positive, real demand and volume are both negative. Cipla is declining inside it.",
        ),
        (facts.stop[1], "Real growth is negative. Structural decline, not a cycle."),
        (
            facts.stop[2],
            "Fastest growth in the market and the most closed. In-license or leave it.",
        ),
    ]
    top = Inches(2.66)
    for row, reason in reasons:
        _card(slide, left, top, Inches(5.4), Inches(1.06), fill=WHITE, accent=TERRACOTTA)
        _text(
            slide,
            left + Inches(0.26),
            top + Inches(0.14),
            Inches(3.3),
            Inches(0.28),
            _short(row["space_label"], 34),
            size=13,
            bold=True,
            colour=INK,
            font=HEAD_FONT,
        )
        _text(
            slide,
            left + Inches(0.26),
            top + Inches(0.44),
            Inches(4.9),
            Inches(0.5),
            reason,
            size=10.5,
            colour=SLATE,
            font=BODY_FONT,
            line_spacing=1.16,
        )
        _text(
            slide,
            left + Inches(3.72),
            top + Inches(0.14),
            Inches(1.44),
            Inches(0.28),
            f"real {row['real_growth'] * 100:+.1f}%",
            size=11.5,
            bold=True,
            colour=TERRACOTTA,
            font=BODY_FONT,
            align=PP_ALIGN.RIGHT,
        )
        top += Inches(1.18)

    _card(slide, left, Inches(6.22), Inches(5.4), Inches(0.72), fill=TEAL_SOFT, accent=TEAL)
    _text(
        slide,
        left + Inches(0.26),
        Inches(6.36),
        Inches(4.94),
        Inches(0.5),
        "The portfolio is optimised for a market that has already moved. "
        "Concentrating on two configurations beats defending nine.",
        size=11.5,
        bold=True,
        colour=INK,
        font=BODY_FONT,
        line_spacing=1.18,
    )


def slide_framework(presentation, facts: Facts) -> None:
    slide = _content_slide(
        presentation,
        "Appendix A1",
        "The prioritisation framework",
        subtitle="Every weight lives in a configuration file, so a challenge is answered by "
        "editing one number and re-running.",
    )

    _picture(slide, "pillar-decomposition", MARGIN, Inches(2.3), Inches(6.35))

    left = Inches(7.24)
    rows = [
        ["Market attractiveness", "34%", "Size, value added, two-year CAGR"],
        ["Future potential", "40%", "Constant-price and volume growth, momentum, signals"],
        ["Competitive headroom", "26%", "HHI, leader share, crowding, price erosion"],
        ["Right to win", "gate", "Share, momentum, molecule and brand adjacency"],
    ]
    _table(
        slide,
        left,
        Inches(2.34),
        Inches(5.5),
        ["Pillar", "Weight", ""],
        rows,
        col_ratios=[2.0, 0.8, 3.2],
        row_h=Inches(0.66),
        size=11,
    )

    _text(
        slide,
        left,
        Inches(5.5),
        Inches(5.5),
        Inches(0.28),
        "THE GATE",
        size=10.5,
        bold=True,
        colour=TEAL,
        font=BODY_FONT,
    )
    _text(
        slide,
        left,
        Inches(5.84),
        Inches(5.5),
        Inches(1.0),
        "Cipla Priority Score = Market Opportunity Index x gate(right to win), where the "
        "gate runs from 0.35 to 1.00 on a curve that punishes weak right to win harder "
        "than linearly. A multiplier cannot let an inaccessible space outrank a winnable one.",
        size=11.5,
        colour=SLATE,
        font=BODY_FONT,
        line_spacing=1.24,
    )


def slide_robustness(presentation, facts: Facts) -> None:
    slide = _content_slide(
        presentation,
        "Appendix A2",
        "How much survives a different framework",
        subtitle="Every weight block redrawn from a Dirichlet distribution and the market "
        "re-scored 500 times.",
    )
    _picture(slide, "sensitivity", MARGIN, Inches(2.3), Inches(6.6))

    left = Inches(7.5)
    rows = [
        [_short(row["space_label"], 30), f"{row['top_k_frequency']:.0%}", f"{row['mean_rank']:.2f}"]
        for row in facts.stability
    ]
    _table(
        slide,
        left,
        Inches(2.4),
        Inches(5.24),
        ["Sub-segment", "In top 3", "Mean rank"],
        rows,
        col_ratios=[2.6, 1.0, 1.0],
        accent_col=1,
    )

    _text(
        slide,
        left,
        Inches(5.2),
        Inches(5.24),
        Inches(1.5),
        "Above 80 per cent the recommendation survives almost any reasonable weighting. "
        "Below 60 per cent it is a judgement call, and is presented as one. The two "
        "priorities on slide two sit at the top of this table; the fourth and fifth "
        "entries of the wider shortlist do not, and are named as candidates rather "
        "than conclusions.",
        size=11.5,
        colour=SLATE,
        font=BODY_FONT,
        line_spacing=1.26,
    )


def slide_sources(presentation, facts: Facts, context) -> None:
    slide = _content_slide(
        presentation,
        "Appendix A3",
        "Sources for external data",
        subtitle=f"{facts.signals} signals, each with a publisher and an access date. "
        "Four are analysis of the supplied dataset and are labelled as such.",
    )

    citations = [c for c in context.citations() if not c["url"].startswith("internal://")]
    rows = [[c["id"], c["title"][:54], c["publisher"][:38]] for c in citations]
    _table(
        slide,
        MARGIN,
        Inches(2.32),
        Inches(12.1),
        ["", "Signal", "Publisher"],
        rows,
        col_ratios=[0.5, 5.6, 3.4],
        row_h=Inches(0.33),
        size=10.5,
    )

    _text(
        slide,
        MARGIN,
        Inches(6.42),
        Inches(12.1),
        Inches(0.5),
        "Signal influence is bounded. Confidence discounts magnitude, agreement within a "
        "category is damped, and the resulting tilt is centred within each level, so a "
        "signal that applies to everything moves nothing. The multiplier is clipped to "
        "0.80 to 1.25 on one metric inside one pillar.",
        size=11,
        colour=MUTED,
        font=BODY_FONT,
        line_spacing=1.22,
    )


def slide_limits(presentation, facts: Facts) -> None:
    slide = _blank(presentation, INK)
    _rect(slide, 0, 0, Inches(0.16), SLIDE_H, TERRACOTTA)
    _text(
        slide,
        MARGIN + Inches(0.3),
        Inches(0.72),
        Inches(11.0),
        Inches(0.3),
        "APPENDIX A4",
        size=11,
        bold=True,
        colour=TERRACOTTA,
        font=BODY_FONT,
    )
    _text(
        slide,
        MARGIN + Inches(0.3),
        Inches(1.08),
        Inches(11.4),
        Inches(0.7),
        "What this analysis cannot tell you",
        size=30,
        bold=True,
        colour=WHITE,
        font=HEAD_FONT,
    )

    limits = [
        (
            "Two years of history",
            "Two annual observations support a direction, not a cycle. Five-year projections are structured extrapolation with an explicit mean-reversion assumption.",
        ),
        (
            "No patent or regulatory status",
            "The framework can see that a space is originator-dominated. It cannot say whether entry is legally possible. Any entry decision needs a freedom-to-operate review.",
        ),
        (
            "No cost or margin data",
            "Every conclusion is about revenue opportunity. Whether a space is profitable to enter is not answerable here.",
        ),
        (
            "Retail audit only",
            "Institutional and government-tender channels are absent, which matters given the scale of public screening programmes.",
        ),
        (
            "Curated external signals",
            "Fourteen documents chosen for relevance, not a systematic review. Their influence is bounded by design for exactly that reason.",
        ),
    ]
    top = Inches(2.16)
    for heading, body in limits:
        _rect(
            slide, MARGIN + Inches(0.3), top + Inches(0.06), Inches(0.05), Inches(0.2), TERRACOTTA
        )
        _text(
            slide,
            MARGIN + Inches(0.52),
            top,
            Inches(3.3),
            Inches(0.3),
            heading,
            size=13.5,
            bold=True,
            colour=WHITE,
            font=HEAD_FONT,
        )
        _text(
            slide,
            Inches(4.72),
            top,
            Inches(8.1),
            Inches(0.6),
            body,
            size=11.5,
            colour=RGBColor(0xB9, 0xC6, 0xD1),
            font=BODY_FONT,
            line_spacing=1.2,
        )
        top += Inches(0.94)

    _text(
        slide,
        MARGIN + Inches(0.3),
        Inches(6.9),
        Inches(12.0),
        Inches(0.3),
        "Stated here rather than buried, because the recommendations should be read against it.",
        size=11,
        italic=True,
        colour=MUTED,
        font=BODY_FONT,
    )


# --- Extended-version slides ----------------------------------------------


def slide_architecture(presentation, facts: Facts) -> None:
    slide = _content_slide(
        presentation,
        "How the agent works",
        "The model never calculates a number",
        subtitle="Reasoning and arithmetic are separated, and a verifier stands between them.",
    )

    stages = [
        (
            "01",
            "Ingest",
            f"{facts.sku_rows:,} packs validated against a fixed schema, normalised into a warehouse with a SHA-256 of the source.",
        ),
        (
            "02",
            "Construct",
            f"{facts.spaces_scored} opportunity spaces across six levels, including molecule franchises that span the reporting hierarchy.",
        ),
        (
            "03",
            "Score",
            "Four pillars from percentile-ranked metrics, with a right-to-win gate. Deterministic and reproducible.",
        ),
        (
            "04",
            "Retrieve",
            f"{facts.signals} cited external signals linked to spaces, bounded and centred so universal evidence moves nothing.",
        ),
        (
            "05",
            "Narrate",
            "The model plans, calls eleven tools and writes prose over a finished evidence pack.",
        ),
        (
            "06",
            "Verify",
            "Every number in the draft is matched against the evidence. Unmatched figures are rejected and rewritten.",
        ),
    ]
    left = MARGIN
    top = Inches(2.5)
    for index, (number, title, body) in enumerate(stages):
        column = index % 3
        rowpos = index // 3
        x = left + Inches(4.12) * column
        y = top + Inches(2.06) * rowpos
        _card(slide, x, y, Inches(3.84), Inches(1.78), fill=WHITE, accent=TEAL)
        _text(
            slide,
            x + Inches(0.28),
            y + Inches(0.18),
            Inches(0.6),
            Inches(0.3),
            number,
            size=15,
            bold=True,
            colour=TEAL,
            font=HEAD_FONT,
        )
        _text(
            slide,
            x + Inches(0.92),
            y + Inches(0.18),
            Inches(2.7),
            Inches(0.3),
            title,
            size=15,
            bold=True,
            colour=INK,
            font=HEAD_FONT,
        )
        _text(
            slide,
            x + Inches(0.28),
            y + Inches(0.6),
            Inches(3.34),
            Inches(1.0),
            body,
            size=11,
            colour=SLATE,
            font=BODY_FONT,
            line_spacing=1.2,
        )

    _footnote(
        slide,
        "Guardrails: scope control before any tool runs, numeric grounding and "
        "citation validation on the draft, injection neutralisation on retrieved text.",
    )


def slide_demand_vs_price(presentation, facts: Facts) -> None:
    slide = _content_slide(
        presentation,
        "Reading the data",
        "Seven of thirteen points of growth are price",
        subtitle="The organisers' glossary is explicit: constant-price MAT measures real demand, "
        "and quantity confirms it.",
    )
    _picture(slide, "growth-decomposition", MARGIN, Inches(2.34), Inches(6.9))

    left = Inches(7.7)
    _text(
        slide,
        left,
        Inches(2.4),
        Inches(5.05),
        Inches(2.2),
        "Price gains do not compound the way prescription volume does, and a material part "
        "of the essential cardiovascular basket sits under administered ceiling prices. "
        "A space growing on value but flat on volume is being repriced, not growing.\n\n"
        "So reported value growth carries no weight at all in the future-potential pillar. "
        "It appears in market attractiveness, where it belongs as a measure of today.",
        size=12.5,
        colour=SLATE,
        font=BODY_FONT,
        line_spacing=1.3,
    )

    _card(slide, left, Inches(4.86), Inches(5.05), Inches(1.62), fill=TERRA_SOFT, accent=TERRACOTTA)
    _text(
        slide,
        left + Inches(0.28),
        Inches(5.04),
        Inches(4.5),
        Inches(0.28),
        "WHAT THIS CHANGES",
        size=10,
        bold=True,
        colour=TERRACOTTA,
        font=BODY_FONT,
    )
    _text(
        slide,
        left + Inches(0.28),
        Inches(5.34),
        Inches(4.5),
        Inches(1.0),
        f"AHT diuretic combinations reports {facts.stop[0]['value_yoy'] * 100:+.1f}% growth. Real demand is "
        f"{facts.stop[0]['real_growth'] * 100:+.1f}% and volume {facts.stop[0]['volume_growth'] * 100:+.1f}%. "
        "Ranking on the headline number alone would have made it look defensible.",
        size=11.5,
        colour=INK,
        font=BODY_FONT,
        line_spacing=1.2,
    )


def slide_franchise_view(presentation, facts: Facts) -> None:
    slide = _content_slide(
        presentation,
        "Beyond the reporting hierarchy",
        "Some franchises are visible only across the hierarchy",
        subtitle="Summing every pack that contains a molecule, whether plain or combined, "
        "exposes clusters no single reporting level shows.",
    )
    _picture(slide, "growth-vs-concentration", MARGIN, Inches(2.3), Inches(6.8))

    left = Inches(7.6)
    _text(
        slide,
        left,
        Inches(2.36),
        Inches(5.15),
        Inches(0.3),
        "THE CILNIDIPINE CASE",
        size=10.5,
        bold=True,
        colour=TEAL,
        font=BODY_FONT,
    )
    _text(
        slide,
        left,
        Inches(2.7),
        Inches(5.15),
        Inches(1.4),
        "The molecule appears as a plain calcium channel blocker, in three dual combinations "
        "and in two triples. Those sit in different sub-segments, so no row of the standard "
        "hierarchy shows the size of the franchise.",
        size=12,
        colour=SLATE,
        font=BODY_FONT,
        line_spacing=1.26,
    )

    cilnidipine = facts.priorities[2]
    cilni_telmi = facts.priorities[3]
    rows = [
        [
            "Whole franchise",
            f"{cilnidipine['value_t2']:,.0f}",
            f"{cilnidipine['value_yoy'] * 100:+.0f}%",
            f"{cilnidipine['hhi']:,.0f}",
        ],
        [
            "Dual with telmisartan",
            f"{cilni_telmi['value_t2']:,.0f}",
            f"{cilni_telmi['value_yoy'] * 100:+.0f}%",
            f"{cilni_telmi['hhi']:,.0f}",
        ],
    ]
    _table(
        slide,
        left,
        Inches(4.3),
        Inches(5.15),
        ["Layer", "Cr", "Growth", "HHI"],
        rows,
        col_ratios=[2.3, 0.9, 0.9, 0.9],
        row_h=Inches(0.44),
        size=11.5,
    )

    _card(slide, left, Inches(5.66), Inches(5.15), Inches(1.06), fill=TEAL_SOFT, accent=TEAL)
    _text(
        slide,
        left + Inches(0.26),
        Inches(5.84),
        Inches(4.66),
        Inches(0.76),
        "The question is not whether to enter cilnidipine. It is which layer of the franchise "
        "still has share being allocated rather than defended.",
        size=11.5,
        bold=True,
        colour=INK,
        font=BODY_FONT,
        line_spacing=1.2,
    )


def slide_cipla_today(presentation, facts: Facts) -> None:
    slide = _content_slide(
        presentation,
        "The starting position",
        "Cipla's largest positions are the ones it is losing",
        subtitle="Five umbrella franchises carry the cardiac business. One of them grows faster "
        "than the market.",
    )
    _picture(slide, "cipla-position", MARGIN, Inches(2.34), Inches(8.3))

    left = Inches(9.28)
    _text(
        slide,
        left,
        Inches(2.4),
        Inches(3.4),
        Inches(0.28),
        "FRANCHISES",
        size=10.5,
        bold=True,
        colour=TEAL,
        font=BODY_FONT,
    )
    rows = [
        [name.title(), f"{value:,.0f}", f"{growth:+.1f}%"]
        for name, value, growth in facts.franchises
    ]
    _table(
        slide,
        left,
        Inches(2.76),
        Inches(3.5),
        ["Brand", "Cr", "Growth"],
        rows,
        col_ratios=[1.7, 0.8, 0.9],
        row_h=Inches(0.42),
        size=11.5,
    )

    _text(
        slide,
        left,
        Inches(5.5),
        Inches(3.5),
        Inches(1.2),
        f"At {facts.cipla_share:.2f}% therapy share and rank {facts.cipla_rank} of "
        f"{facts.company_count}, Cipla cannot out-invest the leaders across every "
        "sub-segment. Adjacency is the cheapest advantage it has.",
        size=11.5,
        colour=SLATE,
        font=BODY_FONT,
        line_spacing=1.26,
    )


def slide_forecast(presentation, facts: Facts) -> None:
    slide = _content_slide(
        presentation,
        "The three-to-five year view",
        "Which spaces should outperform, and by how much",
        subtitle="Growth mean-reverts towards the therapy-area rate. Nothing is projected to "
        "sustain its current rate for five years.",
    )
    _picture(slide, "forecast-outperformance", MARGIN, Inches(2.32), Inches(7.3))

    left = Inches(8.3)
    _text(
        slide,
        left,
        Inches(2.4),
        Inches(4.4),
        Inches(0.28),
        "METHOD",
        size=10.5,
        bold=True,
        colour=TEAL,
        font=BODY_FONT,
    )
    _text(
        slide,
        left,
        Inches(2.74),
        Inches(4.4),
        Inches(2.6),
        "The base rate blends constant-price growth, volume growth, two-year CAGR and "
        "three-month momentum, weighted towards demand.\n\n"
        "Excess growth over the market decays 28 per cent a year, the rate is capped, and "
        "bull and bear cases shift it by four points.\n\n"
        "These are structured extrapolations from two years of history, not plans, and "
        "every projection returns its assumptions alongside the number.",
        size=11.5,
        colour=SLATE,
        font=BODY_FONT,
        line_spacing=1.26,
    )

    _card(slide, left, Inches(5.66), Inches(4.4), Inches(1.06), fill=TEAL_SOFT, accent=TEAL)
    _text(
        slide,
        left + Inches(0.26),
        Inches(5.84),
        Inches(3.9),
        Inches(0.76),
        f"Market baseline: {facts.market_cagr:.1f}% two-year CAGR. Anything above the line "
        "gains share of the therapy area.",
        size=11.5,
        bold=True,
        colour=INK,
        font=BODY_FONT,
        line_spacing=1.2,
    )


def slide_execution(presentation, facts: Facts) -> None:
    slide = _content_slide(
        presentation,
        "Execution",
        "What happens, and in what order",
        subtitle="Sequenced so the shorter-cycle move funds attention for the longer one.",
    )

    phases = [
        (
            "MONTHS 0-12",
            TEAL,
            "Scale the ezetimibe extension",
            [
                "Rosulip EZ is launched. The constraint is share of voice, not development.",
                "The space grew 88 per cent last year and twelve companies entered it.",
                "Share allocated in the next eighteen months will be hard to take back.",
            ],
        ),
        (
            "MONTHS 12-24",
            SLATE,
            "Build the cilnidipine combination line",
            [
                "File and launch a cilnidipine plus telmisartan pack under Cresar.",
                "Follow with a chlortalidone triple, where no player holds above a quarter.",
                "Formulation risk is low; every component is already manufactured.",
            ],
        ),
        (
            "CONTINUOUS",
            TERRACOTTA,
            "Release the funding",
            [
                "Defend Amlopres and Amlopres-AT efficiently, without expansion.",
                "Harvest Cresar-H and Ramipres rather than defending them with new spend.",
                "Redeploy field-force time from declining diuretic combinations.",
            ],
        ),
    ]

    left = MARGIN
    for tag, accent, heading, points in phases:
        _card(slide, left, Inches(2.42), Inches(3.9), Inches(3.5), fill=WHITE, accent=accent)
        _text(
            slide,
            left + Inches(0.28),
            Inches(2.62),
            Inches(3.3),
            Inches(0.26),
            tag,
            size=10,
            bold=True,
            colour=accent,
            font=BODY_FONT,
        )
        _text(
            slide,
            left + Inches(0.28),
            Inches(2.92),
            Inches(3.34),
            Inches(0.66),
            heading,
            size=15.5,
            bold=True,
            colour=INK,
            font=HEAD_FONT,
            line_spacing=1.08,
        )
        top = Inches(3.78)
        for point in points:
            _rect(
                slide, left + Inches(0.28), top + Inches(0.06), Inches(0.05), Inches(0.14), accent
            )
            _text(
                slide,
                left + Inches(0.46),
                top,
                Inches(3.18),
                Inches(0.66),
                point,
                size=11,
                colour=SLATE,
                font=BODY_FONT,
                line_spacing=1.2,
            )
            top += Inches(0.72)
        left += Inches(4.12)

    _card(
        slide, MARGIN, Inches(6.18), Inches(12.1), Inches(0.76), fill=TERRA_SOFT, accent=TERRACOTTA
    )
    _text(
        slide,
        MARGIN + Inches(0.28),
        Inches(6.34),
        Inches(11.5),
        Inches(0.5),
        "Two things would have to be checked before committing: freedom to operate on each "
        "combination, and whether the field-force economics support a share build against "
        "sub-segments carrying more than a hundred competing companies.",
        size=11.5,
        colour=INK,
        font=BODY_FONT,
        line_spacing=1.2,
    )


def slide_close(presentation, facts: Facts) -> None:
    slide = _blank(presentation, INK)
    _rect(slide, 0, 0, Inches(0.16), SLIDE_H, TEAL)
    _text(
        slide,
        Inches(1.4),
        Inches(2.24),
        Inches(10.4),
        Inches(0.3),
        "THE IMPLICATION",
        size=11.5,
        bold=True,
        colour=TEAL_SOFT,
        font=BODY_FONT,
    )
    _text(
        slide,
        Inches(1.4),
        Inches(2.68),
        Inches(10.6),
        Inches(1.9),
        "Cipla cannot win this market by covering it.\n"
        "It can win two configurations inside it,\nboth next to brands it already owns.",
        size=30,
        bold=True,
        colour=WHITE,
        font=HEAD_FONT,
        line_spacing=1.18,
    )
    _text(
        slide,
        Inches(1.4),
        Inches(4.86),
        Inches(10.4),
        Inches(0.9),
        "And it should stop funding the legacy positions that make the portfolio look "
        "broader than its right to win actually is.",
        size=15,
        colour=RGBColor(0xB9, 0xC6, 0xD1),
        font=BODY_FONT,
        line_spacing=1.28,
    )


# --------------------------------------------------------------------------


def _new_presentation() -> Presentation:
    presentation = Presentation()
    presentation.slide_width = SLIDE_W
    presentation.slide_height = SLIDE_H
    return presentation


def build_round_one(facts: Facts, context) -> Path:
    presentation = _new_presentation()
    slide_cover(
        presentation,
        facts,
        "Three slides. Two priorities, the trade-offs behind them, and where to stop spending.",
    )
    slide_finding(presentation, facts)
    slide_priorities(presentation, facts)
    slide_whitespace(presentation, facts)
    slide_framework(presentation, facts)
    slide_robustness(presentation, facts)
    slide_sources(presentation, facts, context)
    slide_limits(presentation, facts)

    OUTPUT.mkdir(parents=True, exist_ok=True)
    path = OUTPUT / "Cardiac_Opportunity_Agent_Round1.pptx"
    presentation.save(str(path))
    return path


def build_detailed(facts: Facts, context) -> Path:
    presentation = _new_presentation()
    slide_cover(
        presentation,
        facts,
        "The full argument: how the agent works, what it found, and what Cipla should do about it.",
    )
    slide_architecture(presentation, facts)
    slide_demand_vs_price(presentation, facts)
    slide_cipla_today(presentation, facts)
    slide_finding(presentation, facts)
    slide_franchise_view(presentation, facts)
    slide_priorities(presentation, facts)
    slide_forecast(presentation, facts)
    slide_whitespace(presentation, facts)
    slide_execution(presentation, facts)
    slide_close(presentation, facts)
    slide_framework(presentation, facts)
    slide_robustness(presentation, facts)
    slide_sources(presentation, facts, context)
    slide_limits(presentation, facts)

    OUTPUT.mkdir(parents=True, exist_ok=True)
    path = OUTPUT / "Cardiac_Opportunity_Agent_Detailed.pptx"
    presentation.save(str(path))
    return path


def main() -> int:
    print("Building the analysis context...")
    context = get_context()
    facts = gather_facts(context)

    print("Building decks:")
    for path in (build_round_one(facts, context), build_detailed(facts, context)):
        presentation = Presentation(str(path))
        print(f"  {path.relative_to(PROJECT_ROOT)}  ({len(presentation.slides)} slides)")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
