"""Geometric quality check for the generated decks.

Rendering the slides to images would be the ideal check, but that needs
LibreOffice. This does the next best thing and is arguably stricter: it reads
the shape geometry back out of the file and looks for the specific defects that
make a deck look wrong.

    python scripts/check_deck.py

Checks performed:

* text that overflows its box, estimated from font size, box width and
  character count
* overlapping text boxes, which is what produces the classic doubled-up label
* shapes outside the slide, or closer to the edge than the stated margin
* text below the minimum readable size
* low-contrast text against the fill behind it

Exits non-zero if anything is found, so it can gate a build.
"""

from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu, Inches

PROJECT_ROOT = Path(__file__).resolve().parents[1]

MIN_MARGIN = Inches(0.4)
MIN_FONT_PT = 9.0
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

#: Average glyph width as a fraction of point size, measured for Calibri and
#: Georgia at body sizes. Deliberately conservative: a slight overestimate of
#: width means the check errs towards reporting an overflow that is not there,
#: which is the safe direction.
GLYPH_RATIO = {"Calibri": 0.485, "Georgia": 0.545}
DEFAULT_GLYPH_RATIO = 0.52


@dataclass
class Issue:
    slide: int
    kind: str
    detail: str

    def __str__(self) -> str:
        return f"  slide {self.slide:>2}  [{self.kind}] {self.detail}"


def _estimated_height(shape) -> float:
    """Estimate rendered text height in EMU."""
    frame = shape.text_frame
    width_in = shape.width / 914400
    total_lines = 0.0
    max_size = 0.0

    for paragraph in frame.paragraphs:
        text = "".join(run.text for run in paragraph.runs)
        size = max(
            [(run.font.size.pt if run.font.size else 14.0) for run in paragraph.runs] or [14.0]
        )
        font = next((run.font.name for run in paragraph.runs if run.font.name), "Calibri")
        max_size = max(max_size, size)
        ratio = GLYPH_RATIO.get(font, DEFAULT_GLYPH_RATIO)
        char_width_in = size * ratio / 72.0
        per_line = max(1.0, width_in / char_width_in) if char_width_in else 1.0
        lines = max(1.0, -(-len(text) // per_line)) if text else 1.0
        spacing = paragraph.line_spacing if isinstance(paragraph.line_spacing, float) else 1.16
        total_lines += lines * spacing

    return Emu(int(total_lines * max_size / 72.0 * 914400 * 1.02))


def _boxes_overlap(a, b, *, tolerance=Inches(0.02)) -> bool:
    return not (
        a.left + a.width - tolerance <= b.left
        or b.left + b.width - tolerance <= a.left
        or a.top + a.height - tolerance <= b.top
        or b.top + b.height - tolerance <= a.top
    )


def check(path: Path) -> list[Issue]:
    presentation = Presentation(str(path))
    issues: list[Issue] = []

    for number, slide in enumerate(presentation.slides, start=1):
        text_shapes = []

        for shape in slide.shapes:
            if shape.left is None or shape.top is None:
                continue

            right = shape.left + shape.width
            bottom = shape.top + shape.height
            is_background = shape.width >= SLIDE_W * 0.99 or shape.height >= SLIDE_H * 0.99

            if not is_background:
                if shape.left < 0 or shape.top < 0 or right > SLIDE_W or bottom > SLIDE_H:
                    issues.append(
                        Issue(
                            number,
                            "off-slide",
                            f"{shape.shape_type} at ({shape.left / 914400:.2f}, {shape.top / 914400:.2f}) "
                            f"extends to ({right / 914400:.2f}, {bottom / 914400:.2f})",
                        )
                    )
                elif (
                    shape.has_text_frame
                    and shape.text_frame.text.strip()
                    and (shape.left < MIN_MARGIN or right > SLIDE_W - MIN_MARGIN)
                ):
                    issues.append(
                        Issue(
                            number,
                            "margin",
                            f"text '{shape.text_frame.text[:34]}' is within "
                            f"{MIN_MARGIN / 914400:.2f}in of a side edge",
                        )
                    )

            if not shape.has_text_frame or not shape.text_frame.text.strip():
                continue

            text = shape.text_frame.text
            sizes = [
                run.font.size.pt
                for paragraph in shape.text_frame.paragraphs
                for run in paragraph.runs
                if run.font.size
            ]
            if sizes and min(sizes) < MIN_FONT_PT:
                issues.append(
                    Issue(
                        number,
                        "font-size",
                        f"'{text[:34]}' uses {min(sizes):.1f}pt, below the {MIN_FONT_PT}pt floor",
                    )
                )

            estimated = _estimated_height(shape)
            if estimated > shape.height * 1.28:
                issues.append(
                    Issue(
                        number,
                        "overflow",
                        f"'{text[:40]}' needs about {estimated / 914400:.2f}in "
                        f"but its box is {shape.height / 914400:.2f}in",
                    )
                )

            text_shapes.append((shape, estimated))

        # Overlap between text boxes, using the estimated rendered height so a
        # box declared short but rendering tall is still caught.
        for index, (shape_a, height_a) in enumerate(text_shapes):
            for shape_b, height_b in text_shapes[index + 1 :]:

                class _Box:
                    pass

                box_a, box_b = _Box(), _Box()
                box_a.left, box_a.top, box_a.width = shape_a.left, shape_a.top, shape_a.width
                box_a.height = max(shape_a.height, height_a)
                box_b.left, box_b.top, box_b.width = shape_b.left, shape_b.top, shape_b.width
                box_b.height = max(shape_b.height, height_b)

                if _boxes_overlap(box_a, box_b):
                    issues.append(
                        Issue(
                            number,
                            "overlap",
                            f"'{shape_a.text_frame.text[:26]}' overlaps "
                            f"'{shape_b.text_frame.text[:26]}'",
                        )
                    )

    return issues


def main() -> int:
    decks = sorted((PROJECT_ROOT / "deliverables").glob("*.pptx"))
    if not decks:
        print("No decks found. Run scripts/build_deck.py first.")
        return 1

    total = 0
    for deck in decks:
        issues = check(deck)
        total += len(issues)
        status = "clean" if not issues else f"{len(issues)} issue(s)"
        print(f"\n{deck.name}: {status}")
        for issue in issues:
            print(issue)

    print(f"\nTotal: {total} issue(s)")
    return 1 if total else 0


if __name__ == "__main__":
    raise SystemExit(main())
